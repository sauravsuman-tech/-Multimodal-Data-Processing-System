from typing import List, Dict
from pathlib import Path
import frontmatter, markdown
import tempfile
import os
import shutil

from PyPDF2 import PdfReader
import docx
from pptx import Presentation
from PIL import Image
import pytesseract

# audio/video
import whisper
from youtube_dl import YoutubeDL
from pydub import AudioSegment

def parse_pdf(path: Path) -> List[Dict]:
    text = []
    try:
        r = PdfReader(str(path))
        for p in r.pages:
            text.append(p.extract_text() or "")
    except Exception as e:
        text = [f"(pdf parse error) {e}"]
    return [{"text": "\n".join(text), "meta": {"type": "pdf"}}]

def parse_docx(path: Path) -> List[Dict]:
    doc = docx.Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text]
    return [{"text": "\n".join(paragraphs), "meta": {"type": "docx"}}]

def parse_pptx(path: Path) -> List[Dict]:
    prs = Presentation(str(path))
    slides_text = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        slides_text.append("\n".join(texts))
    return [{"text": "\n\n".join(slides_text), "meta": {"type": "pptx"}}]

def parse_txt(path: Path) -> List[Dict]:
    return [{"text": path.read_text(encoding="utf-8", errors="ignore"), "meta": {"type": "txt"}}]

def parse_md(path: Path) -> List[Dict]:
    content = path.read_text(encoding="utf-8", errors="ignore")
    # optional: strip frontmatter
    try:
        fm = frontmatter.loads(content)
        text = fm.content
    except Exception:
        text = content
    plain = markdown.markdown(text)
    return [{"text": plain, "meta": {"type": "md"}}]

def parse_image(path: Path) -> List[Dict]:
    try:
        img = Image.open(str(path))
        txt = pytesseract.image_to_string(img)
    except Exception as e:
        txt = f"(ocr error) {e}"
    return [{"text": txt, "meta": {"type": "image"}}]

_whisper_model = None
def get_whisper():
    global _whisper_model
    if _whisper_model is None:
        _whisper_model = whisper.load_model("small")  # or base
    return _whisper_model

def transcribe_audio_file(path: Path) -> List[Dict]:
    model = get_whisper()
    res = model.transcribe(str(path))
    text = res.get("text", "")
    return [{"text": text, "meta": {"type": "audio"}}]

def download_youtube_audio(url: str, out_path: Path) -> Path:
    ydl_opts = {
        "format": "bestaudio/best",
        "outtmpl": str(out_path),
        "quiet": True,
        "no_warnings": True,
        "postprocessors": [{"key": "FFmpegExtractAudio", "preferredcodec": "mp3", "preferredquality": "192"}]
    }
    with YoutubeDL(ydl_opts) as ydl:
        ydl.download([url])
    return out_path

def parse_media(path: Path) -> List[Dict]:
    tmp_wav = path.with_suffix(".wav")
    try:
        audio = AudioSegment.from_file(str(path))
        audio.export(str(tmp_wav), format="wav")
        docs = transcribe_audio_file(tmp_wav)
    except Exception as e:
        docs = [{"text": f"(audio parse error) {e}", "meta": {"type": "audio"}}]
    finally:
        if tmp_wav.exists():
            tmp_wav.unlink()
    return docs

def parse_file(pathlike) -> List[Dict]:
    p = Path(pathlike)
    suf = p.suffix.lower()
    if suf == ".pdf":
        return parse_pdf(p)
    if suf in [".docx"]:
        return parse_docx(p)
    if suf in [".pptx"]:
        return parse_pptx(p)
    if suf in [".txt"]:
        return parse_txt(p)
    if suf in [".md", ".markdown"]:
        return parse_md(p)
    if suf in [".png", ".jpg", ".jpeg", ".tiff", ".bmp"]:
        return parse_image(p)
    if suf in [".mp3", ".wav", ".mp4", ".m4a", ".aac", ".flac"]:
        return parse_media(p)
    # fallback: try reading as text
    try:
        return parse_txt(p)
    except Exception:
        return [{"text": "(unrecognized file type)", "meta": {"type": "unknown"}}]
